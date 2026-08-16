"""24-slide faculty-grade dark HPC briefing. Measured numbers only."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# Palette
NAVY = RGBColor(0x05, 0x08, 0x14)
NAVY2 = RGBColor(0x0A, 0x12, 0x20)
PANEL = RGBColor(0x0E, 0x16, 0x24)
PANEL2 = RGBColor(0x12, 0x1C, 0x2E)
LINE = RGBColor(0x1C, 0x2C, 0x42)
GRID = RGBColor(0x12, 0x1C, 0x2E)
WHITE = RGBColor(0xF4, 0xF7, 0xFB)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CORAL = RGBColor(0xFB, 0x71, 0x85)
ORB_C = RGBColor(0x0A, 0x2A, 0x38)
ORB_V = RGBColor(0x18, 0x12, 0x32)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Segoe UI"
N = 24

P100 = [7.589, 5.663, 4.497, 4.535, 3.899, 4.327]
P500 = [43.732, 33.735, 24.414, 20.455, 17.947, 17.287]
WORKERS = ["1", "2", "4", "6", "8", "12"]


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def run(p, text, size, color, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT
    return r


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def tb(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run(p, text, size, color, bold)
    return box


def rect(slide, l, t, w, h, color, oval=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL if oval else MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill(sh, color)
    return sh


def rrect(slide, l, t, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, color)
    if line:
        stroke(sh, line, 1.0)
    else:
        sh.line.fill.background()
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def layered_bg(slide):
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, Inches(8.6), Inches(-1.4), Inches(6.2), Inches(4.4), ORB_C, oval=True)
    rect(slide, Inches(-1.8), Inches(4.6), Inches(5.4), Inches(4.2), ORB_V, oval=True)
    # faint grid
    for i in range(1, 16):
        y = Inches(0.42 * i)
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Emu(6350))
        fill(ln, GRID)
    for i in range(1, 20):
        x = Inches(0.67 * i)
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, 0, Emu(6350), H)
        fill(ln, GRID)
    # rails
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), H)
    fill(rail, CYAN)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.03))
    fill(top, CYAN)
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.18), W, Inches(0.32))
    fill(bot, NAVY2)


def footer(slide, n: int, section: str):
    tb(slide, Inches(0.42), Inches(7.18), Inches(8.6), Inches(0.32), f"CSE 471  ·  Parallel Log Intelligence  ·  {section}", 11, MUTED, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    tb(slide, Inches(10.4), Inches(7.18), Inches(2.5), Inches(0.32), f"{n}  /  {N}", 11, CYAN, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def kicker(slide, text: str, color=VIOLET):
    tb(slide, Inches(0.48), Inches(0.18), Inches(12.3), Inches(0.28), text.upper(), 11, color, True)


def heading(slide, text: str, y=0.42):
    tb(slide, Inches(0.48), Inches(y), Inches(12.3), Inches(0.58), text, 30, WHITE, True)


def sub(slide, text: str, y=0.98):
    tb(slide, Inches(0.48), Inches(y), Inches(12.3), Inches(0.38), text, 15, MUTED)


def ask_rail(slide):
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), H)
    fill(rail, CYAN)


def badge(slide, label="ASK THE ROOM"):
    sh = rrect(slide, Inches(0.48), Inches(0.18), Inches(2.55), Inches(0.34), RGBColor(0x0A, 0x2C, 0x36), CYAN)
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], label, 10, CYAN, True)


def chip(slide, l, t, w, h, text, color):
    sh = rrect(slide, l, t, w, h, RGBColor(0x0A, 0x2C, 0x36), color)
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text, 12, color, True)
    return sh


def bullets(slide, l, t, w, h, items, size=15, color=WHITE):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run(p, "▸  " + item, size, color)
    return box


def _solid_fill_xml(hex_color: str):
    spPr = OxmlElement("c:spPr")
    solid = OxmlElement("a:solidFill")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", hex_color)
    solid.append(srgb)
    spPr.append(solid)
    return spPr


def style_chart(chart):
    chart.has_legend = False
    ser = chart.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = AMBER
    ser.format.line.fill.background()
    try:
        chart.font.color.rgb = MUTED
        chart.font.size = Pt(10)
        chart.font.name = FONT
    except Exception:
        pass
    for axis in (chart.category_axis, chart.value_axis):
        axis.has_major_gridlines = True
        axis.tick_labels.font.color.rgb = MUTED
        axis.tick_labels.font.size = Pt(10)
        axis.tick_labels.font.name = FONT
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.paragraphs[0].clear()
    run(chart.value_axis.axis_title.text_frame.paragraphs[0], "Seconds (mean)", 10, MUTED)
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.paragraphs[0].clear()
    run(chart.category_axis.axis_title.text_frame.paragraphs[0], "Workers p", 10, MUTED)
    # dark plot + chart space
    plot = chart._element.chart.plotArea
    plot.append(_solid_fill_xml("0E1624"))
    sp = OxmlElement("c:spPr")
    solid = OxmlElement("a:solidFill")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", "0A1220")
    solid.append(srgb)
    sp.append(solid)
    chart._element.append(sp)


def add_column_chart(slide, l, t, w, h, values, series_name):
    data = CategoryChartData()
    data.categories = WORKERS
    data.add_series(series_name, values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, data)
    style_chart(gf.chart)
    return gf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    layered_bg(s)
    return s


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Title
    s = new_slide(prs)
    kicker(s, "CSE 471  ·  High-Performance Computing  ·  Stage 2", VIOLET)
    tb(s, Inches(0.48), Inches(1.35), Inches(12.2), Inches(1.1), "Parallel Log Intelligence", 40, WHITE, True)
    tb(s, Inches(0.48), Inches(2.45), Inches(12.2), Inches(0.7), "We split a large log across CPU cores, merge one correct report, and optionally explain the evidence with local AI.", 18, MUTED)
    chip(s, Inches(0.48), Inches(3.45), Inches(2.35), Inches(0.4), "HPC engine", CYAN)
    chip(s, Inches(3.0), Inches(3.45), Inches(2.55), Inches(0.4), "Honest metrics", AMBER)
    chip(s, Inches(5.72), Inches(3.45), Inches(2.85), Inches(0.4), "Optional local AI", TEAL)
    chip(s, Inches(8.75), Inches(3.45), Inches(2.2), Inches(0.4), "No GPU path", VIOLET)
    rrect(s, Inches(0.48), Inches(4.25), Inches(12.3), Inches(1.55), PANEL, LINE)
    tb(s, Inches(0.7), Inches(4.42), Inches(12.0), Inches(0.35), "Hardware measured", 12, MUTED, True)
    tb(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(0.7), "Intel i5-1235U  ·  12 logical processors  ·  12 GB RAM  ·  Iris Xe (no CUDA device)  ·  ~15 min faculty briefing", 16, WHITE)
    footer(s, 1, "Purpose")
    notes(s, "Thank the room. One sentence purpose. Hardware is a suitability decision: Iris Xe cannot run CUDA, so this is multi-paradigm CPU HPC.")

    # 2 Ask
    s = new_slide(prs)
    ask_rail(s)
    badge(s)
    tb(s, Inches(0.48), Inches(1.1), Inches(12.3), Inches(2.3), "When something breaks in a product you run — an app, a bank portal, a campus system — where do you look first?", 28, WHITE, True)
    tb(s, Inches(0.48), Inches(3.6), Inches(12.3), Inches(0.45), "Gut feel is fine. Ten seconds. Two or three answers.", 16, CYAN)
    rrect(s, Inches(0.48), Inches(4.4), Inches(12.3), Inches(1.4), PANEL, LINE)
    tb(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(0.35), "Hint — not the answer yet", 12, MUTED, True)
    tb(s, Inches(0.7), Inches(5.05), Inches(12.0), Inches(0.5), "Logs: timestamped records of logins, errors, slow requests, and suspicion.", 16, WHITE)
    footer(s, 2, "Purpose")
    notes(s, "Pause. Collect answers. Then: the first question is rarely rewrite the product. It is what do the logs say?")

    # 3 Purpose
    s = new_slide(prs)
    kicker(s, "Main purpose", CYAN)
    heading(s, "One correct report. Many cores.")
    sub(s, "The engine does HPC. The product shows it. AI only translates evidence.")
    steps = [
        ("01", "Decompose", "Byte chunks, full lines only"),
        ("02", "Analyze", "Workers parse in parallel"),
        ("03", "Merge", "Associative partial results"),
        ("04", "Explain", "Optional local AI on totals"),
    ]
    for i, (n, t, b) in enumerate(steps):
        x = 0.48 + i * 3.2
        rrect(s, Inches(x), Inches(1.6), Inches(3.0), Inches(3.55), PANEL, LINE)
        tb(s, Inches(x + 0.2), Inches(1.85), Inches(2.6), Inches(0.4), n, 14, CYAN, True)
        tb(s, Inches(x + 0.2), Inches(2.4), Inches(2.6), Inches(0.7), t, 22, WHITE, True)
        tb(s, Inches(x + 0.2), Inches(3.25), Inches(2.6), Inches(1.4), b, 15, MUTED)
        if i < 3:
            tb(s, Inches(x + 2.85), Inches(3.0), Inches(0.4), Inches(0.4), "→", 18, AMBER, True, PP_ALIGN.CENTER)
    footer(s, 3, "Purpose")
    notes(s, "This is the whole product in four verbs. AI is optional and never reads the raw multi-hundred-MB file.")

    # 4 Problem
    s = new_slide(prs)
    kicker(s, "The business problem", CORAL)
    heading(s, "The cost is not disk space")
    sub(s, "It is time-to-understand.")
    bullets(s, Inches(0.48), Inches(1.5), Inches(12.3), Inches(2.0), [
        "Logs are the black box of software — logins, failures, slow APIs, suspicious IPs.",
        "Delayed reading means delayed recovery, delayed customer message, delayed trust.",
        "One day’s file can be hundreds of megabytes; one core wastes the rest of the CPU.",
    ])
    labels = [("Incident", CYAN), ("Confusion", AMBER), ("Slow decision", CORAL)]
    for i, (lab, col) in enumerate(labels):
        x = 0.48 + i * 4.2
        rrect(s, Inches(x), Inches(4.0), Inches(3.9), Inches(1.7), PANEL, col)
        tb(s, Inches(x), Inches(4.5), Inches(3.9), Inches(0.7), lab, 22, col, True, PP_ALIGN.CENTER)
        if i < 2:
            tb(s, Inches(x + 3.7), Inches(4.55), Inches(0.5), Inches(0.5), "→", 20, WHITE, True, PP_ALIGN.CENTER)
    footer(s, 4, "Purpose")
    notes(s, "Observability platforms exist for this. Today is the same idea as a measurable parallel-computing prototype — not a commercial SIEM.")

    # 5 Sequential KPIs
    s = new_slide(prs)
    kicker(s, "Measured pain  ·  i5-1235U", AMBER)
    heading(s, "Sequential is already slow")
    kpis = [
        ("100 MB", "~7.6 s", "one core  ·  901,610 lines"),
        ("500 MB", "~44 s", "one core  ·  4.5M lines"),
        ("Read-only", "~0.05 s", "100 MB bytes only — not the bottleneck"),
    ]
    cols = [AMBER, AMBER, TEAL]
    for i, ((a, b, c), col) in enumerate(zip(kpis, cols)):
        x = 0.48 + i * 4.2
        rrect(s, Inches(x), Inches(1.55), Inches(3.95), Inches(3.7), PANEL, LINE)
        tb(s, Inches(x + 0.25), Inches(1.85), Inches(3.45), Inches(0.4), a, 14, MUTED, True)
        tb(s, Inches(x + 0.25), Inches(2.45), Inches(3.45), Inches(1.1), b, 36, col, True)
        tb(s, Inches(x + 0.25), Inches(3.8), Inches(3.45), Inches(1.0), c, 14, WHITE)
    tb(s, Inches(0.48), Inches(5.5), Inches(12.3), Inches(0.7), "Interactive investigation cannot wait. Real production logs are often larger.", 15, MUTED)
    footer(s, 5, "Engine")
    notes(s, "Speak only measured numbers. Full tables in PERFORMANCE.md. Do not round into marketing.")

    # 6 Ask analogy
    s = new_slide(prs)
    ask_rail(s)
    badge(s)
    tb(s, Inches(0.48), Inches(0.7), Inches(12.3), Inches(1.5), "A 500-page incident report: one reader for every page, or eight readers of chapters who then combine counts?", 22, WHITE, True)
    rrect(s, Inches(0.48), Inches(2.45), Inches(5.9), Inches(3.4), PANEL, CORAL)
    tb(s, Inches(0.7), Inches(2.7), Inches(5.5), Inches(0.35), "SERIAL", 12, CORAL, True)
    tb(s, Inches(0.7), Inches(3.2), Inches(5.5), Inches(0.8), "1 reader  ·  500 pages  ·  slow", 22, WHITE, True)
    tb(s, Inches(0.7), Inches(4.2), Inches(5.5), Inches(1.1), "Correct, but the rest of the team sits idle.", 15, MUTED)
    rrect(s, Inches(6.9), Inches(2.45), Inches(5.9), Inches(3.4), PANEL, TEAL)
    tb(s, Inches(7.12), Inches(2.7), Inches(5.5), Inches(0.35), "PARALLEL", 12, TEAL, True)
    tb(s, Inches(7.12), Inches(3.2), Inches(5.5), Inches(0.8), "8 readers  ·  chapters  ·  one report", 20, WHITE, True)
    chip(s, Inches(7.12), Inches(4.5), Inches(5.45), Inches(0.7), "Only if nobody double-counts pages", AMBER)
    footer(s, 6, "Engine")
    notes(s, "That is HPC. Same math. More workers. One correct answer if chunk boundaries and the merge are honest.")

    # 7 Pipeline
    s = new_slide(prs)
    kicker(s, "How it works", CYAN)
    heading(s, "Split. Analyze. Merge.")
    pipe = [
        ("File", "Offline log"),
        ("Chunks", "Newline aligned"),
        ("Workers", "Independent map"),
        ("Partials", "Counts only"),
        ("Reduce", "Deterministic"),
        ("Report", "Findings + top-N"),
    ]
    for i, (t, b) in enumerate(pipe):
        x = 0.4 + i * 2.15
        rrect(s, Inches(x), Inches(1.55), Inches(2.0), Inches(2.55), PANEL, CYAN if i < 5 else TEAL)
        tb(s, Inches(x + 0.1), Inches(1.8), Inches(1.8), Inches(0.35), f"{i+1:02d}", 12, CYAN, True, PP_ALIGN.CENTER)
        tb(s, Inches(x + 0.08), Inches(2.25), Inches(1.84), Inches(0.7), t, 16, WHITE, True, PP_ALIGN.CENTER)
        tb(s, Inches(x + 0.08), Inches(3.05), Inches(1.84), Inches(0.7), b, 12, MUTED, False, PP_ALIGN.CENTER)
        if i < 5:
            tb(s, Inches(x + 1.88), Inches(2.45), Inches(0.3), Inches(0.4), "→", 14, AMBER, True, PP_ALIGN.CENTER)
    rrect(s, Inches(0.48), Inches(4.4), Inches(12.3), Inches(1.5), PANEL, TEAL)
    tb(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(0.35), "Optional local AI", 13, TEAL, True)
    tb(s, Inches(0.7), Inches(5.05), Inches(12.0), Inches(0.55), "Explains the report — aggregates and findings only. Never the raw file. If Ollama is down, analytics still work.", 15, WHITE)
    footer(s, 7, "Engine")
    notes(s, "Faculty: domain decomposition, map, associative reduce. Business: many cores, one scoreboard. Newline align so no line is split or duplicated.")

    # 8 System A/B
    s = new_slide(prs)
    kicker(s, "Architecture freeze", VIOLET)
    heading(s, "Engine and product stay separate")
    rrect(s, Inches(0.48), Inches(1.5), Inches(6.05), Inches(4.35), PANEL, CYAN)
    tb(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.35), "SYSTEM A  ·  HPC ENGINE", 13, CYAN, True)
    bullets(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(3.3), [
        "CLI-runnable without FastAPI or React",
        "Chunk → parse → histograms → merge",
        "Backends: process, dynamic, OpenMP, MPI",
        "Faculty can grade HPC without the website",
    ], 15)
    rrect(s, Inches(6.8), Inches(1.5), Inches(6.05), Inches(4.35), PANEL, VIOLET)
    tb(s, Inches(7.02), Inches(1.7), Inches(5.6), Inches(0.35), "SYSTEM B  ·  PRODUCT", 13, VIOLET, True)
    bullets(s, Inches(7.02), Inches(2.2), Inches(5.6), Inches(3.3), [
        "React + FastAPI + SQLite + JWT",
        "Jobs: queued → running → aggregating → completed",
        "Dashboards, benchmarks, optional Ollama",
        "Does not contain the core parallel algorithm",
    ], 15)
    footer(s, 8, "Engine")
    notes(s, "If faculty asks for HPC, open the CLI. If industry asks for a product, show the UI and job lifecycle.")

    # 9 Backends
    s = new_slide(prs)
    kicker(s, "Multi-paradigm CPU HPC", CYAN)
    heading(s, "One job. Four ways to run it.")
    grid = [
        ("ProcessPool", CYAN, "Several Python processes. Default. Strong on 100–500 MB files."),
        ("Dynamic", AMBER, "Many small chunks. Teaching load-balance. Helps uneven work."),
        ("OpenMP", TEAL, "Native C, shared-memory threads. Often fastest on small files."),
        ("MPI", VIOLET, "mpi4py + mpiexec on one node. Message-passing model; not a cluster yet."),
    ]
    for i, (t, col, b) in enumerate(grid):
        r, c = divmod(i, 2)
        x, y = 0.48 + c * 6.4, 1.45 + r * 2.3
        rrect(s, Inches(x), Inches(y), Inches(6.15), Inches(2.1), PANEL, col)
        tb(s, Inches(x + 0.25), Inches(y + 0.25), Inches(5.7), Inches(0.45), t, 20, col, True)
        tb(s, Inches(x + 0.25), Inches(y + 0.85), Inches(5.7), Inches(0.95), b, 15, WHITE)
    footer(s, 9, "Engine")
    notes(s, "Do not deep-dive APIs. Product chooses how to run; engine does the work. GPU acceleration is out of scope on Iris Xe.")

    # 10 CLO map
    s = new_slide(prs)
    kicker(s, "Why this project for CSE 471", VIOLET)
    heading(s, "Course outcomes, not a demo toy")
    clos = [
        ("CLO 1", "Measure & benchmark", "Speedup Sp, efficiency Ep, strong + weak scaling, backend compare."),
        ("CLO 2", "CLI / UNIX-style path", "hpc_engine.analyze, WSL2, submit / run / collect scripts."),
        ("CLO 3", "Jobs & environments", "queued→completed; cluster-lite scripts; venv as module analogue."),
    ]
    for i, (a, b, c) in enumerate(clos):
        x = 0.48 + i * 4.2
        rrect(s, Inches(x), Inches(1.5), Inches(4.0), Inches(4.35), PANEL, VIOLET)
        tb(s, Inches(x + 0.22), Inches(1.75), Inches(3.55), Inches(0.4), a, 13, VIOLET, True)
        tb(s, Inches(x + 0.22), Inches(2.3), Inches(3.55), Inches(1.0), b, 20, WHITE, True)
        tb(s, Inches(x + 0.22), Inches(3.5), Inches(3.55), Inches(1.9), c, 14, MUTED)
    footer(s, 10, "Course")
    notes(s, "Point at CLO1 as the faculty hook: we publish timed matrices, not vibes. CLO3 is honest — not Slurm on a university cluster.")

    # 11 Week map
    s = new_slide(prs)
    kicker(s, "Week-by-week coverage", VIOLET)
    heading(s, "The course map on one slide")
    rows = [
        ("W2–4", "Models + metrics", "Four backends; Sp, Ep; domain decomposition"),
        ("W5–7", "OpenMP", "Native C worker, #pragma omp parallel for"),
        ("W8–10", "MPI", "Single-node mpi4py + mpiexec — teaching path"),
        ("W11–14", "GPU / CUDA weeks", "Evaluated; unsuitable on Iris Xe — documented"),
    ]
    for i, (w, t, b) in enumerate(rows):
        y = 1.48 + i * 1.2
        rrect(s, Inches(0.48), Inches(y), Inches(12.3), Inches(1.08), PANEL, LINE)
        tb(s, Inches(0.7), Inches(y + 0.28), Inches(1.7), Inches(0.5), w, 16, CYAN, True)
        tb(s, Inches(2.5), Inches(y + 0.28), Inches(3.3), Inches(0.5), t, 16, WHITE, True)
        tb(s, Inches(6.0), Inches(y + 0.28), Inches(6.5), Inches(0.5), b, 15, MUTED)
    footer(s, 11, "Course")
    notes(s, "Say the GPU weeks out loud: objective (f) is evaluate suitability. Choosing CUDA would fail on this laptop.")

    # 12 Why logs
    s = new_slide(prs)
    kicker(s, "Why this application", AMBER)
    heading(s, "Why logs — not a textbook multiply")
    items = [
        ("Industry workload", "Observability, SRE, security operations — faculty from industry recognize the domain."),
        ("Embarrassingly parallel after chunking", "Workers do not talk until reduce. A clean parallel story."),
        ("Correctness is non-trivial", "Bad boundaries duplicate or drop lines. Parity tests prove we are not faster-but-wrong."),
        ("Measurable on student hardware", "100–500 MB on an i5 is enough for T1, Sp, Ep without a cluster."),
    ]
    for i, (t, b) in enumerate(items):
        y = 1.45 + i * 1.2
        rrect(s, Inches(0.48), Inches(y), Inches(12.3), Inches(1.08), PANEL, LINE)
        tb(s, Inches(0.7), Inches(y + 0.15), Inches(12.0), Inches(0.32), t, 15, AMBER, True)
        tb(s, Inches(0.7), Inches(y + 0.5), Inches(12.0), Inches(0.45), b, 14, WHITE)
    footer(s, 12, "Course")
    notes(s, "This answers why not matrix multiply: real workload, hard correctness, honest laptop-scale evidence.")

    # 13 Protocol
    s = new_slide(prs)
    kicker(s, "Benchmark protocol", AMBER)
    heading(s, "How we measure — so you can trust it")
    proto = [
        ("Warm-up", "Discarded. Then timed runs."),
        ("Clock", "time.perf_counter, same file/parser/analysis."),
        ("Formulas", "Sp = T1 / Tp     Ep = Sp / p"),
        ("Scaling", "Strong (fixed size) + weak (~50 MB/worker)."),
        ("Honesty", "No psutil CPU%. No invented times."),
        ("Machine", "i5-1235U · AC power · Best Performance."),
    ]
    for i, (t, b) in enumerate(proto):
        r, c = divmod(i, 3)
        x, y = 0.48 + c * 4.2, 1.5 + r * 2.35
        rrect(s, Inches(x), Inches(y), Inches(4.0), Inches(2.15), PANEL, AMBER)
        tb(s, Inches(x + 0.22), Inches(y + 0.3), Inches(3.55), Inches(0.4), t, 16, AMBER, True)
        tb(s, Inches(x + 0.22), Inches(y + 0.85), Inches(3.55), Inches(0.95), b, 14, WHITE)
    footer(s, 13, "Evidence")
    notes(s, "Faculty should hear the protocol before the pretty chart. That is scientific HPC, not a marketing graph.")

    # 14 Chart 100MB
    s = new_slide(prs)
    kicker(s, "Experiment A  ·  strong scaling", AMBER)
    heading(s, "100 MB wall-clock vs workers")
    add_column_chart(s, Inches(0.4), Inches(1.35), Inches(8.4), Inches(5.3), P100, "Mean s")
    rrect(s, Inches(8.95), Inches(1.45), Inches(3.9), Inches(5.15), PANEL, AMBER)
    tb(s, Inches(9.15), Inches(1.7), Inches(3.5), Inches(0.35), "BEST", 12, MUTED, True)
    tb(s, Inches(9.15), Inches(2.15), Inches(3.5), Inches(0.7), "p = 8", 28, WHITE, True)
    tb(s, Inches(9.15), Inches(2.95), Inches(3.5), Inches(0.7), "S8 ≈ 1.95×", 22, AMBER, True)
    tb(s, Inches(9.15), Inches(3.75), Inches(3.5), Inches(0.5), "3.899 s mean", 14, WHITE)
    tb(s, Inches(9.15), Inches(4.4), Inches(3.5), Inches(1.6), "p=12 is slower than p=8. Overhead beat remaining work. Expected on hybrid P/E cores.", 13, MUTED)
    footer(s, 14, "Evidence")
    notes(s, "Pause on the chart. Sequential 7.589 s. Best 3.899 s at 8 workers. Efficiency 0.24. Do not say 8×.")

    # 15 Chart 500MB
    s = new_slide(prs)
    kicker(s, "Experiment A  ·  strong scaling", AMBER)
    heading(s, "500 MB wall-clock vs workers")
    add_column_chart(s, Inches(0.4), Inches(1.35), Inches(8.4), Inches(5.3), P500, "Mean s")
    rrect(s, Inches(8.95), Inches(1.45), Inches(3.9), Inches(5.15), PANEL, TEAL)
    tb(s, Inches(9.15), Inches(1.7), Inches(3.5), Inches(0.35), "BEST", 12, MUTED, True)
    tb(s, Inches(9.15), Inches(2.15), Inches(3.5), Inches(0.7), "p = 12", 28, WHITE, True)
    tb(s, Inches(9.15), Inches(2.95), Inches(3.5), Inches(0.7), "S12 ≈ 2.53×", 22, AMBER, True)
    tb(s, Inches(9.15), Inches(3.75), Inches(3.5), Inches(0.5), "17.287 s mean", 14, WHITE)
    tb(s, Inches(9.15), Inches(4.4), Inches(3.5), Inches(1.6), "Larger parse work amortizes spawn. Efficiency still 0.21. Sublinear — and published.", 13, MUTED)
    footer(s, 15, "Evidence")
    notes(s, "Sequential 43.732 s. Best 17.287 s. Throughput about 11.4 MB/s to 28.9 MB/s. Honest 2.53×.")

    # 16 CPU-bound
    s = new_slide(prs)
    kicker(s, "Bottleneck classification", TEAL)
    heading(s, "CPU-bound — not a disk story")
    sub(s, "100 MB sequential profiles. Warm-up discarded.")
    bars = [
        ("Read-only", 0.048, 0.05, MUTED, "0.048 s"),
        ("Parse-only", 5.583, 5.6, CYAN, "5.583 s"),
        ("Parse + analyze", 8.955, 8.96, AMBER, "8.955 s"),
    ]
    maxw = 11.6
    for i, (lab, val, _, col, cap) in enumerate(bars):
        y = 1.7 + i * 1.35
        tb(s, Inches(0.48), Inches(y), Inches(2.5), Inches(0.4), lab, 14, WHITE, True)
        bw = Inches(max(0.25, maxw * (val / 8.955)))
        sh = rrect(s, Inches(3.1), Inches(y), bw, Inches(0.55), col)
        tb(s, Inches(3.1) + bw + Inches(0.15), Inches(y), Inches(2.0), Inches(0.55), cap, 14, col, True)
    tb(s, Inches(0.48), Inches(5.7), Inches(12.3), Inches(0.6), "Disk read is ~0.5% of parse+analyze. Parallel CPUs help because parsing dominates.", 15, MUTED)
    footer(s, 16, "Evidence")
    notes(s, "This answers the I/O question. If it were disk-bound, extra workers would not help. They do — modestly, as the charts show.")

    # 17 Ask honesty
    s = new_slide(prs)
    ask_rail(s)
    badge(s)
    tb(s, Inches(0.48), Inches(0.75), Inches(12.3), Inches(1.6), "Would you rather hear “12 times faster on 12 cores” — or an honest ~2× we can defend?", 24, WHITE, True)
    bullets(s, Inches(0.48), Inches(2.6), Inches(12.3), Inches(3.5), [
        "Hybrid chip: 2 performance cores + 8 efficiency cores + hyper-threading = 12 logical processors.",
        "Extra workers share weaker cores; spawn, pickle, and merge cost time (Amdahl serial fraction).",
        "On 100 MB, 12 workers lost to 8. On 500 MB, 12 still wins wall-clock with Ep ≈ 0.21.",
        "Claiming S12 = 12 on this part would be false. Faculty prefer measured 2×.",
    ], 16)
    footer(s, 17, "Evidence")
    notes(s, "Let them answer. Then: serious rooms prefer a number that survives due diligence.")

    # 18 Benefits
    s = new_slide(prs)
    kicker(s, "Who benefits", TEAL)
    heading(s, "Three audiences. One system.")
    bens = [
        ("Faculty", VIOLET, "Gradable HPC: OpenMP, MPI, ProcessPool, Sp/Ep, parity tests, published protocol."),
        ("Student / CV", CYAN, "Real repo, honest tables, product UI, CLI independence — not a toy notebook."),
        ("Operator", AMBER, "Faster offline investigation on CPUs you already own. Local AI on aggregates, not the raw dump."),
    ]
    for i, (t, col, b) in enumerate(bens):
        x = 0.48 + i * 4.2
        rrect(s, Inches(x), Inches(1.5), Inches(4.0), Inches(4.35), PANEL, col)
        tb(s, Inches(x + 0.22), Inches(1.8), Inches(3.55), Inches(0.6), t, 20, col, True)
        tb(s, Inches(x + 0.22), Inches(2.7), Inches(3.55), Inches(2.6), b, 15, WHITE)
    footer(s, 18, "Industry")
    notes(s, "Keep benefits concrete. Do not invent dollar ROI. Time-to-insight is the business translation of wall-clock.")

    # 19 vs typical tools
    s = new_slide(prs)
    kicker(s, "Industry context", AMBER)
    heading(s, "Typical tools vs this prototype")
    sub(s, "Same problem class. Different job.")
    rrect(s, Inches(0.48), Inches(1.55), Inches(6.05), Inches(4.3), PANEL, MUTED)
    tb(s, Inches(0.7), Inches(1.75), Inches(5.6), Inches(0.4), "SPLUNK / ELASTIC / DATADOG", 13, MUTED, True)
    bullets(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(3.2), [
        "Production ingest, search, alerting, ops",
        "Licensed platforms, long-running clusters",
        "They already parallelize internally",
        "Not what we claim to replace",
    ], 15)
    rrect(s, Inches(6.8), Inches(1.55), Inches(6.05), Inches(4.3), PANEL, CYAN)
    tb(s, Inches(7.02), Inches(1.75), Inches(5.6), Inches(0.4), "THIS PROJECT", 13, CYAN, True)
    bullets(s, Inches(7.02), Inches(2.3), Inches(5.6), Inches(3.2), [
        "Teaches and measures the HPC mechanism",
        "Four backends you can switch and time",
        "Correctness contract + published Sp/Ep",
        "Local, CPU-only, laptop-scale, auditable",
    ], 15)
    footer(s, 19, "Industry")
    notes(s, "Industry uses parallel log processing inside those products. We expose the mechanism so a course — and a hiring manager — can see it.")

    # 20 When it fits
    s = new_slide(prs)
    kicker(s, "When this approach is the right fit", TEAL)
    heading(s, "Use this pattern when…")
    fits = [
        ("CPU-only sites", "No NVIDIA device. Suitability says skip CUDA."),
        ("Offline large files", "Incident dumps, not live Kafka — yet."),
        ("You must prove scaling", "Sp, Ep, strong/weak — not a vendor slide."),
        ("Privacy / local AI", "Aggregates only. Raw logs stay on the machine."),
        ("Need a product path", "Same engine behind CLI and a real UI."),
        ("Need correctness", "Faster is worthless if counts disagree."),
    ]
    for i, (t, b) in enumerate(fits):
        r, c = divmod(i, 3)
        x, y = 0.48 + c * 4.2, 1.5 + r * 2.35
        rrect(s, Inches(x), Inches(y), Inches(4.0), Inches(2.15), PANEL, TEAL)
        tb(s, Inches(x + 0.22), Inches(y + 0.3), Inches(3.55), Inches(0.55), t, 16, TEAL, True)
        tb(s, Inches(x + 0.22), Inches(y + 0.95), Inches(3.55), Inches(0.85), b, 14, WHITE)
    footer(s, 20, "Industry")
    notes(s, "This is the honest sales slide: not always, but when these constraints hold.")

    # 21 Trust
    s = new_slide(prs)
    kicker(s, "Trust pillars", TEAL)
    heading(s, "Faster is worthless if it is wrong")
    pillars = [
        ("Correctness", "Sequential ≡ parallel (pytest, all backends when present)."),
        ("Performance", "Measurable Sp on 100 MB and 500 MB."),
        ("Scalability", "Strong + weak scaling documented."),
        ("Evidence", "JSON summaries. No fabricated times."),
        ("Modularity", "Engine runs with no website."),
        ("Privacy", "Ollama sees aggregates only."),
    ]
    for i, (t, b) in enumerate(pillars):
        r, c = divmod(i, 3)
        x, y = 0.48 + c * 4.2, 1.5 + r * 2.35
        rrect(s, Inches(x), Inches(y), Inches(4.0), Inches(2.15), PANEL, TEAL)
        tb(s, Inches(x + 0.22), Inches(y + 0.3), Inches(3.55), Inches(0.45), t, 16, TEAL, True)
        tb(s, Inches(x + 0.22), Inches(y + 0.9), Inches(3.55), Inches(0.9), b, 14, WHITE)
    footer(s, 21, "Trust")
    notes(s, "Five course pillars plus privacy. Keep the slide human — point at parity tests if challenged.")

    # 22 Scope
    s = new_slide(prs)
    kicker(s, "Scope honesty", CORAL)
    heading(s, "Clear claims. No surprises.")
    rrect(s, Inches(0.48), Inches(1.5), Inches(6.05), Inches(4.35), PANEL, TEAL)
    tb(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.4), "WE DID", 14, TEAL, True)
    bullets(s, Inches(0.7), Inches(2.25), Inches(5.6), Inches(3.3), [
        "Laptop multi-paradigm HPC",
        "ProcessPool, OpenMP, single-node MPI",
        "100 MB and 500 MB measured",
        "Web app + independent CLI",
    ], 15)
    rrect(s, Inches(6.8), Inches(1.5), Inches(6.05), Inches(4.35), PANEL, CORAL)
    tb(s, Inches(7.02), Inches(1.7), Inches(5.6), Inches(0.4), "WE DID NOT", 14, CORAL, True)
    bullets(s, Inches(7.02), Inches(2.25), Inches(5.6), Inches(3.3), [
        "GPU acceleration on this hardware",
        "Multi-node / Slurm cluster",
        "Real-time streaming ingest",
        "A production SIEM product",
    ], 15)
    footer(s, 22, "Trust")
    notes(s, "Honesty is a feature. Stage 2 is course-aligned CPU HPC. Stages 3–6 remain future.")

    # 23 Close
    s = new_slide(prs)
    kicker(s, "Close", CYAN)
    heading(s, "Many cores. One correct report.")
    bullets(s, Inches(0.48), Inches(1.5), Inches(12.3), Inches(1.8), [
        "Split large logs across CPUs you already have.",
        "Merge evidence you can audit — sequential equals parallel.",
        "Optional local AI on the summary, not the raw file.",
    ], 18)
    rrect(s, Inches(0.48), Inches(3.6), Inches(12.3), Inches(1.5), PANEL, AMBER)
    tb(s, Inches(0.7), Inches(3.95), Inches(12.0), Inches(0.8), "~1.95× on 100 MB     ·     ~2.53× on 500 MB", 26, AMBER, True, PP_ALIGN.CENTER)
    tb(s, Inches(0.48), Inches(5.4), Inches(12.3), Inches(0.5), "github.com/galibhasan720/Parallel-Log-Analysis", 18, CYAN, False, PP_ALIGN.CENTER)
    footer(s, 23, "Close")
    notes(s, "One breath. Invite questions from faculty and from business in the same room.")

    # 24 Q&A
    s = new_slide(prs)
    ask_rail(s)
    badge(s, "ASK THE ROOM  ·  Q&A")
    qs = [
        ("Business", CYAN, "Where does investigation time cost you the most today?"),
        ("Faculty", VIOLET, "Which backend would you stress-test first — processes, OpenMP, or MPI?"),
        ("Anyone", AMBER, "What would you refuse to believe without a measured table?"),
    ]
    for i, (who, col, q) in enumerate(qs):
        y = 0.85 + i * 1.45
        rrect(s, Inches(0.48), Inches(y), Inches(12.3), Inches(1.28), PANEL, col)
        tb(s, Inches(0.7), Inches(y + 0.15), Inches(2.2), Inches(0.3), who.upper(), 12, col, True)
        tb(s, Inches(0.7), Inches(y + 0.52), Inches(11.8), Inches(0.55), q, 18, WHITE, True)
    tb(s, Inches(0.48), Inches(5.35), Inches(12.3), Inches(0.85), "Thank you", 32, WHITE, True, PP_ALIGN.CENTER)
    footer(s, 24, "Close")
    notes(s, "Let them pick a question. If you lack a number, say you will check PERFORMANCE.md. Never invent speedup.")

    out = Path(__file__).resolve().parents[1] / "docs" / "Parallel_Log_Intelligence_Faculty.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print(build())
